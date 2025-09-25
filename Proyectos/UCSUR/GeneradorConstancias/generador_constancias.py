from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
import pandas as pd
import os
import comtypes.client

# Configuración de rutas
excel_path = r"C:\Users\manue\Downloads\OneDrive_2024-11-29\Generador de Constancias\Participantes.xlsx"
pptx_template_path = r"C:\Users\manue\Downloads\OneDrive_2024-11-29\Generador de Constancias\Plantilla.pptx"
output_dir = r"C:\Users\manue\Downloads\OneDrive_2024-11-29\Generador de Constancias\pdfs"
pdf_output_dir = os.path.join(output_dir, "PDFs")

# Crear los directorios de salida si no existen
os.makedirs(output_dir, exist_ok=True)
os.makedirs(pdf_output_dir, exist_ok=True)

# Cargar archivo Excel
participants = pd.read_excel(excel_path, sheet_name="Lista")
names = participants["Nombres"].dropna()

# Color personalizado (RGB del color extraído de la plantilla)
highlight_color = RGBColor(34, 70, 144)  # Azul extraído de la plantilla

# Inicializar PowerPoint una vez
powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
powerpoint.Visible = 1

# Función para convertir PPTX a PDF
def convert_to_pdf(powerpoint, pptx_path, pdf_path):
    try:
        presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
        presentation.SaveAs(pdf_path, 32)  # 32 es el formato para PDF
        presentation.Close()
    except Exception as e:
        print(f"Error al convertir {pptx_path} a PDF: {e}")

# Generar presentaciones personalizadas
pptx_paths = []
for name in names:
    presentation = Presentation(pptx_template_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and "Nombres" in shape.text:
                shape.text = shape.text.replace("Nombres", name)
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Gotham Black"
                        run.font.bold = True
                        run.font.size = Pt(24)
                        run.font.color.rgb = highlight_color  # Aplicar color

    # Guardar presentación personalizada
    pptx_path = os.path.join(output_dir, f"{name}.pptx")
    presentation.save(pptx_path)
    pptx_paths.append(pptx_path)

# Convertir presentaciones a PDF
for pptx_path in pptx_paths:
    pdf_name = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
    pdf_path = os.path.join(pdf_output_dir, pdf_name)
    convert_to_pdf(powerpoint, pptx_path, pdf_path)

# Cerrar PowerPoint
powerpoint.Quit()

print(f"Presentaciones y PDFs generados correctamente en: {pdf_output_dir}")